import os
import traceback
import fitz  # pymupdf
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from config import SUPPORTED_EXTENSIONS, RECURSIVE_SCAN


def get_file_extension(file_path):
    """
    获取文件后缀，统一转成小写。
    例如：
    demo.PDF -> .pdf
    """
    return os.path.splitext(file_path)[1].lower()


def scan_documents(docs_dir, recursive=True):
    """
    扫描文档目录，返回所有支持格式的文件路径。
    """
    file_paths = []
    if recursive:
        for root, dirs, files in os.walk(docs_dir):
            for filename in files:
                file_path = os.path.join(root, filename)
                ext = get_file_extension(file_path)
                if ext in SUPPORTED_EXTENSIONS:
                    file_paths.append(file_path)
    else:
        for filename in os.listdir(docs_dir):
            file_path = os.path.join(docs_dir, filename)
            if not os.path.isfile(file_path):
                continue
            ext = get_file_extension(file_path)
            if ext in SUPPORTED_EXTENSIONS:
                file_paths.append(file_path)
    return file_paths


def dataframe_to_text(df):
    """
    将 DataFrame 转成适合 embedding 的文本。
    """
    df = df.fillna("")
    texts = []
    # 表头
    columns = [str(col) for col in df.columns]
    texts.append("表头：" + " | ".join(columns))
    # 每一行转成一段文本
    for idx, row in df.iterrows():
        row_items = []
        for col in df.columns:
            value = row[col]
            row_items.append(f"{col}: {value}")
        texts.append("；".join(row_items))
    return "\n".join(texts)


def parse_txt(file_path):
    """
    解析 txt / md 文件。
    """
    encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312"]
    last_error = None
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read()
            return text
        except UnicodeDecodeError as e:
            last_error = e
    raise last_error


def parse_pdf(file_path):
    """
    解析 PDF。
    适合有文字层的 PDF。
    如果是扫描版 PDF，需要 OCR。
    """
    texts = []
    doc = fitz.open(file_path)
    for page_index, page in enumerate(doc):
        page_text = page.get_text("text")
        if page_text.strip():
            texts.append(f"[第{page_index + 1}页]\n{page_text}")
    doc.close()
    return "\n\n".join(texts)


def parse_docx(file_path):
    """
    解析 docx 文档。
    提取段落和表格内容。
    """
    doc = Document(file_path)
    texts = []
    # 提取段落
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    # 提取表格
    for table_index, table in enumerate(doc.tables):
        texts.append(f"[表格{table_index + 1}]")
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                row_texts.append(cell_text)
            texts.append(" | ".join(row_texts))
    return "\n".join(texts)


def parse_csv(file_path):
    """
    解析 CSV。
    将表格转换成文本。
    """
    encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312"]
    last_error = None
    for enc in encodings:
        try:
            df = pd.read_csv(file_path, encoding=enc)
            return dataframe_to_text(df)
        except UnicodeDecodeError as e:
            last_error = e
    raise last_error


def parse_excel(file_path):
    """
    解析 xlsx / xls。
    读取所有 sheet。
    """
    excel_data = pd.read_excel(file_path, sheet_name=None)
    texts = []
    for sheet_name, df in excel_data.items():
        texts.append(f"[Sheet: {sheet_name}]")
        texts.append(dataframe_to_text(df))
    return "\n\n".join(texts)


def parse_pptx(file_path):
    """
    解析 pptx。
    提取每页幻灯片中的文本框内容。
    """
    prs = Presentation(file_path)
    texts = []
    for slide_index, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            texts.append(f"[第{slide_index + 1}页幻灯片]")
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def parse_html(file_path):
    """
    解析 HTML。
    去掉 script/style，只提取可见文本。
    """
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.extract()
    text = soup.get_text(separator="\n")
    lines = []
    for line in text.splitlines():
        line = line.strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def parse_document(file_path):
    """
    根据文件扩展名，自动选择解析函数。
    """
    ext = get_file_extension(file_path)
    if ext in [".txt", ".md"]:
        text = parse_txt(file_path)
    elif ext == ".pdf":
        text = parse_pdf(file_path)
    elif ext == ".docx":
        text = parse_docx(file_path)
    elif ext == ".csv":
        text = parse_csv(file_path)
    elif ext in [".xlsx", ".xls"]:
        text = parse_excel(file_path)
    elif ext in [".html", ".htm"]:
        text = parse_html(file_path)
    elif ext == ".pptx":
        text = parse_pptx(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {ext}")
    return {
        "source": file_path,
        "extension": ext,
        "text": text
    }


def load_all_documents(docs_dir, recursive=RECURSIVE_SCAN):
    """
    加载 docs_dir 下所有支持的文档。
    """
    file_paths = scan_documents(docs_dir, recursive=recursive)
    print(f"发现支持的文档数量: {len(file_paths)}")
    documents = []
    for file_path in file_paths:
        print(f"\n正在解析文档: {file_path}")
        try:
            doc = parse_document(file_path)
            if doc["text"].strip():
                documents.append(doc)
                print(f"解析成功，文本长度: {len(doc['text'])}")
            else:
                print("警告：解析结果为空，已跳过。")
        except Exception as e:
            print(f"解析失败: {file_path}")
            print(f"错误信息: {e}")
            traceback.print_exc()
    return documents

# --coding:utf-8--
